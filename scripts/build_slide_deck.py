"""
build_slide_deck.py — generate FlightCast_Pitch_Deck.pptx from screenshots.

Builds a 12-slide pitch deck (plus 3 backup Q&A slides) styled to match the
dashboard's dark theme. Embeds dashboard screenshots from
``screenshot/`` plus the architecture diagram from ``web/public/``.

Run from project root:
    python scripts/build_slide_deck.py

Output: docs/FlightCast_Pitch_Deck.pptx

If a screenshot is missing, the script renders a styled placeholder
rectangle so the slide still has the right shape — paste your image
into the placeholder in PowerPoint and delete the rectangle.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── Paths ──────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "screenshot"
OUT = ROOT / "docs" / "FlightCast_Pitch_Deck.pptx"
ARCH = ROOT / "web" / "public" / "architecture.png"

IMG = {
    "comparison": SHOTS / "Screenshot 2026-05-10 175052.png",
    "dashboard_hero": SHOTS / "Screenshot 2026-05-10 175117.png",
    "time_travel": SHOTS / "Screenshot 2026-05-10 175132.png",
    "rainbow": SHOTS / "Screenshot 2026-05-10 175147.png",
    "drift_tiles": SHOTS / "Screenshot 2026-05-10 175158.png",
    "math": SHOTS / "Screenshot 2026-05-10 175215.png",
    "architecture": ARCH,
}

# ─── Theme (matches the dashboard) ──────────────────────────────────
BG = RGBColor(0x0A, 0x0F, 0x1F)
BG_CARD = RGBColor(0x11, 0x18, 0x2B)
TEXT_PRIMARY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_TERTIARY = RGBColor(0x64, 0x74, 0x8B)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_SOFT = RGBColor(0xC4, 0xB5, 0xFD)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x34, 0xD3, 0x99)

FONT = "Inter"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12  # main slides; backups don't affect page numbers


# ─── Helpers ────────────────────────────────────────────────────────
def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    color=TEXT_PRIMARY,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_bullets(
    slide, left, top, width, height, items, *, size=16, color=TEXT_PRIMARY, bullet=VIOLET
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        dot = p.add_run()
        dot.text = "•  "
        dot.font.size = Pt(size + 4)
        dot.font.color.rgb = bullet
        dot.font.name = FONT
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = FONT
    return box


def add_image(slide, key, left, top, width, height, *, label=None):
    """Embed image from IMG[key]; render placeholder if file missing."""
    p = IMG.get(key)
    if p and Path(p).exists():
        slide.shapes.add_picture(str(p), left, top, width, height)
        return
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_CARD
    rect.line.color.rgb = TEXT_TERTIARY
    rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = f"[ paste {label or key} screenshot here ]"
    r.font.size = Pt(14)
    r.font.color.rgb = TEXT_TERTIARY
    r.font.italic = True
    r.font.name = FONT


def add_pill(slide, left, top, width, height, text, *, fg=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.5
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_CARD
    rect.line.color.rgb = fg
    rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.color.rgb = fg
    r.font.bold = True
    r.font.name = FONT


def add_footer(slide, n):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        "FlightCast  ·  MariaDB Hackathon Malaysia 2026  ·  Innovation Track",
        size=10,
        color=TEXT_TERTIARY,
    )
    add_text(
        slide,
        Inches(11.6),
        Inches(7.05),
        Inches(1.2),
        Inches(0.3),
        f"{n} / {TOTAL}",
        size=10,
        color=TEXT_TERTIARY,
        align=PP_ALIGN.RIGHT,
    )


def add_eyebrow(slide, text, color=AMBER):
    add_text(
        slide,
        Inches(0.7),
        Inches(0.55),
        Inches(12),
        Inches(0.35),
        text,
        size=12,
        color=color,
        bold=True,
    )


def add_title(slide, text, size=42):
    add_text(
        slide,
        Inches(0.7),
        Inches(0.95),
        Inches(12),
        Inches(1.0),
        text,
        size=size,
        color=TEXT_PRIMARY,
        bold=True,
    )


def add_subtitle(slide, text):
    add_text(
        slide,
        Inches(0.7),
        Inches(1.95),
        Inches(12),
        Inches(0.6),
        text,
        size=18,
        color=TEXT_SECONDARY,
    )


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout


# ─── Slides ─────────────────────────────────────────────────────────


def slide_01_title(prs):
    s = new_slide(prs)
    add_bg(s)
    # Decorative top bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VIOLET
    bar.line.fill.background()
    # Title and subtitle, centered vertically
    add_text(
        s,
        Inches(1),
        Inches(2.6),
        Inches(11.33),
        Inches(1.6),
        "FlightCast",
        size=92,
        color=TEXT_PRIMARY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(4.1),
        Inches(11.33),
        Inches(0.7),
        "Self-auditing ML predictions on MariaDB temporal tables",
        size=22,
        color=VIOLET_SOFT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(5.9),
        Inches(11.33),
        Inches(0.4),
        "Low Yan Cheng  ·  TP070056  ·  APU Malaysia",
        size=15,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(6.35),
        Inches(11.33),
        Inches(0.4),
        "MariaDB Hackathon Malaysia 2026  ·  Innovation Track",
        size=12,
        color=TEXT_TERTIARY,
        align=PP_ALIGN.CENTER,
    )
    notes(
        s,
        "FlightCast — self-auditing machine learning predictions, built on MariaDB "
        "temporal tables. By Low Yan Cheng, APU Malaysia, for the Innovation Track.",
    )


def slide_02_problem(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 1  ·  THE PROBLEM", color=AMBER)
    add_title(s, "What did your AI predict last month?")
    add_subtitle(s, "Most ML systems can't answer that question.")
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(3.5),
        [
            "AI models retrain every week and overwrite yesterday's predictions.",
            "When auditors, regulators, or post-mortems ask what the model said three months ago — most teams cannot answer.",
            "The industry workaround bolts on MLflow + Weights & Biases + a custom audit log. Three services. Fragile and expensive.",
            "The audit gap is structural, not a tooling problem.",
        ],
        bullet=AMBER,
    )
    add_footer(s, 2)
    notes(
        s,
        "Production AI retrains every week. Each retrain overwrites the previous "
        "predictions. So three months later, when an auditor asks 'what did your "
        "model predict on January 15?', most companies cannot answer. The industry "
        "workaround is to bolt on three external services. FlightCast solves it "
        "inside the database.",
    )


def slide_03_insight(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 1  ·  THE INSIGHT", color=VIOLET_SOFT)
    add_title(s, "Five SQL keywords only MariaDB has.")
    add_subtitle(s, "FOR  SYSTEM_TIME  AS  OF")
    # Big code chip
    code = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(2.9),
        Inches(12),
        Inches(1.1),
    )
    code.adjustments[0] = 0.15
    code.fill.solid()
    code.fill.fore_color.rgb = BG_CARD
    code.line.color.rgb = VIOLET
    code.line.width = Pt(0.75)
    tf = code.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "SELECT * FROM forecasts FOR SYSTEM_TIME AS OF '2026-01-15';"
    r.font.size = Pt(20)
    r.font.color.rgb = VIOLET_SOFT
    r.font.name = "Consolas"
    # Comparison row of pills
    pill_y = Inches(4.6)
    add_pill(s, Inches(0.7), pill_y, Inches(2.8), Inches(0.55), "MariaDB 11.8  ✓ NATIVE", fg=EMERALD)
    add_pill(s, Inches(3.7), pill_y, Inches(2.8), Inches(0.55), "MySQL  ✗ SYNTAX ERROR", fg=AMBER)
    add_pill(s, Inches(6.7), pill_y, Inches(2.8), Inches(0.55), "PostgreSQL  ⚠ EXTENSION", fg=AMBER)
    add_pill(s, Inches(9.7), pill_y, Inches(3.0), Inches(0.55), "SQLite  ✗ NOT SUPPORTED", fg=AMBER)
    add_text(
        s,
        Inches(0.7),
        Inches(5.55),
        Inches(12),
        Inches(1.2),
        "MariaDB is the only mainstream database that ships system-versioned "
        "tables out of the box. One keyword unlocks the entire audit story.",
        size=16,
        color=TEXT_SECONDARY,
    )
    add_footer(s, 3)
    notes(
        s,
        "MariaDB has a feature called system-versioned tables. Five English-like "
        "words — FOR SYSTEM TIME AS OF a date — let you query any past state of "
        "any row. MySQL gives a syntax error. PostgreSQL needs an add-on. SQLite "
        "has nothing. MariaDB is the only mainstream database that ships this "
        "out of the box.",
    )


def slide_04_solution(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 2  ·  THE SOLUTION", color=VIOLET_SOFT)
    add_title(s, "The audit trail IS the database.")
    add_subtitle(s, "Every prediction the model ever made — queryable forever.")
    add_image(s, "dashboard_hero", Inches(0.7), Inches(2.85), Inches(8.5), Inches(4.0), label="dashboard hero")
    # Bullet column
    add_bullets(
        s,
        Inches(9.5),
        Inches(2.9),
        Inches(3.4),
        Inches(4.0),
        [
            "No MLflow tracking server.",
            "No Weights & Biases.",
            "No application audit log.",
            "Just MariaDB system-versioned tables + MAPIE conformal prediction.",
        ],
        size=14,
        bullet=VIOLET,
    )
    add_footer(s, 4)
    notes(
        s,
        "The audit trail isn't a service we add later — it's structural. Every "
        "prediction is system-versioned at INSERT time. We layer a math result "
        "called MAPIE conformal prediction on top to guarantee 90% coverage. "
        "Two ideas combined: temporal SQL plus conformal math.",
    )


def slide_05_time_travel(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 2  ·  DEMO  ·  TIME TRAVEL", color=VIOLET_SOFT)
    add_title(s, "Drag the slider. Travel through model versions.")
    add_subtitle(s, "Each tick is a real ROW_START timestamp from a committed transaction.")
    add_image(s, "time_travel", Inches(0.7), Inches(2.85), Inches(8.5), Inches(4.0), label="time-travel chart")
    add_bullets(
        s,
        Inches(9.5),
        Inches(2.9),
        Inches(3.4),
        Inches(4.0),
        [
            "6 committed model versions live in MariaDB.",
            "Slider position = real database timestamp.",
            "Chart re-queries the database on every drag.",
            "Empirical coverage card jumps with the slider — calibrated vs drift in one number.",
        ],
        size=14,
        bullet=VIOLET,
    )
    add_footer(s, 5)
    notes(
        s,
        "Drag the slider and the dashboard re-queries the database at that exact "
        "micro-second. Six committed model versions, all live, all queryable "
        "through one SQL statement. No log replay, no shadow tables, no external "
        "service.",
    )


def slide_06_rainbow(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 2  ·  DEMO  ·  ALL HISTORY", color=VIOLET_SOFT)
    add_title(s, "Six model versions, one SQL query.")
    add_subtitle(s, "FOR  SYSTEM_TIME  ALL  returns every version of every row.")
    add_image(s, "rainbow", Inches(0.7), Inches(2.85), Inches(8.5), Inches(4.0), label="rainbow chart")
    add_bullets(
        s,
        Inches(9.5),
        Inches(2.9),
        Inches(3.4),
        Inches(4.0),
        [
            "One query → 6 versions × 30 days = 180 rows.",
            "MLflow workflow: 6 separate API calls + Python join.",
            "Bands stacked tightly = healthy.",
            "Bands diverging = drift in real time.",
        ],
        size=14,
        bullet=VIOLET,
    )
    add_footer(s, 6)
    notes(
        s,
        "Six committed model versions, each in a different colour, each with its "
        "own confidence band, all on one chart. Behind it: a single MariaDB "
        "query. In a normal MLflow setup you would issue six API calls and "
        "stitch them together in Python.",
    )


def slide_07_drift(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 2  ·  THE PUNCHLINE", color=AMBER)
    add_title(s, "Drift caught by ONE SQL query.")
    add_subtitle(s, "The database itself is the auditor.")
    # Drift tiles screenshot, large and centered
    add_image(s, "drift_tiles", Inches(0.7), Inches(3.0), Inches(12), Inches(2.0), label="drift tiles")
    # Big numbers row underneath
    add_text(
        s,
        Inches(0.7),
        Inches(5.4),
        Inches(12),
        Inches(0.5),
        "One GROUP BY against FOR SYSTEM_TIME ALL.  No MLflow.  No tracking server.  No data engineer.",
        size=14,
        color=TEXT_SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    # Three accent stats
    add_text(s, Inches(0.7), Inches(6.0), Inches(4), Inches(0.5),
             "30+ pp", size=24, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.5), Inches(4), Inches(0.4),
             "coverage drop", size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.7), Inches(6.0), Inches(4), Inches(0.5),
             "3.3 ×", size=24, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.7), Inches(6.5), Inches(4), Inches(0.4),
             "Winkler escalation", size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.7), Inches(6.0), Inches(4), Inches(0.5),
             "1 query", size=24, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.7), Inches(6.5), Inches(4), Inches(0.4),
             "to detect both", size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    add_footer(s, 7)
    notes(
        s,
        "This is the punchline. We trained six batches. The first four use normal "
        "noise — the world is calm. Batches five and six simulate a fuel-price "
        "shock, with much louder noise. Result: the first four hit 91.2% "
        "coverage, dead on target. The last two collapse to 58.7%. That is more "
        "than thirty percentage points off. The database itself caught it, with "
        "one SQL query, in real time.",
    )


def slide_08_math(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 3  ·  THE MATH LAYER", color=VIOLET_SOFT)
    add_title(s, "A guarantee, not a claim.")
    add_subtitle(s, "MAPIE conformal prediction — published 2005, peer-reviewed.")
    add_image(s, "math", Inches(0.7), Inches(2.85), Inches(8.5), Inches(3.5), label="conformal math")
    add_bullets(
        s,
        Inches(9.5),
        Inches(2.9),
        Inches(3.4),
        Inches(4.0),
        [
            "Vovk 2005, Lei 2018 — peer-reviewed coverage theorem.",
            "Used in finance, insurance, and clinical trials.",
            "We didn't invent it. We applied it on top of MariaDB so the math and the audit trail share one storage layer.",
        ],
        size=14,
        bullet=VIOLET,
    )
    add_footer(s, 8)
    notes(
        s,
        "The 90% coverage isn't marketing. It's a published theorem from 2005, "
        "used by major banks and insurance firms. We didn't invent it. We "
        "applied it on top of MariaDB so the math and the audit trail live in "
        "the same place.",
    )


def slide_09_architecture(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 3  ·  ARCHITECTURE", color=VIOLET_SOFT)
    add_title(s, "Five layers. Zero external dependencies.")
    add_subtitle(s, "The audit trail is structural, not bolted on.")
    add_image(s, "architecture", Inches(1.5), Inches(2.7), Inches(10.3), Inches(4.0), label="architecture diagram")
    add_text(
        s,
        Inches(0.7),
        Inches(6.65),
        Inches(12),
        Inches(0.4),
        "Notably absent:  MLflow  ·  Weights & Biases  ·  DataDog  ·  Evidently AI",
        size=12,
        color=TEXT_TERTIARY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, 9)
    notes(
        s,
        "Five layers. Notably absent from this diagram: MLflow, Weights and "
        "Biases, DataDog, Evidently. The forecasts table is system-versioned at "
        "INSERT time, so the audit trail is structural — not a service you bolt "
        "on later.",
    )


def slide_10_comparison(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 3  ·  DIFFERENTIATION", color=VIOLET_SOFT)
    add_title(s, "The moat is the math layer, not the slider.")
    add_subtitle(s, "What FlightCast replaces in the standard MLOps stack.")
    add_image(s, "comparison", Inches(0.7), Inches(2.85), Inches(12), Inches(3.6), label="comparison table")
    add_text(
        s,
        Inches(0.7),
        Inches(6.55),
        Inches(12),
        Inches(0.4),
        "Every advantage on the right is a problem the database itself solved.",
        size=14,
        color=VIOLET_SOFT,
        align=PP_ALIGN.CENTER,
        italic=True,
    )
    add_footer(s, 10)
    notes(
        s,
        "Compared to the standard stack: external tracking server, replaced by "
        "atomic INSERT-time versioning. Custom drift dashboards, replaced by "
        "one SQL query. Replay infrastructure, replaced by FOR SYSTEM TIME. "
        "The moat is the math layer, not the slider.",
    )


def slide_11_open(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "ACT 4  ·  OPEN INFRASTRUCTURE", color=EMERALD)
    add_title(s, "MIT-licensed.  Fork it.  Build on it.")
    add_subtitle(s, "Public infrastructure for the MariaDB ecosystem.")
    # Three-column layout
    col_w = Inches(3.9)
    col_y = Inches(3.0)
    col_h = Inches(3.4)
    gap = Inches(0.25)
    x = Inches(0.7)
    cards = [
        ("5 hero SQL queries", "All MariaDB-only.\nNone run on MySQL or PostgreSQL.\nReusable as-is.", VIOLET),
        ("3 open issues", "Per-tier MAPIE recalibration.\nReal-data prototypes.\nApache Airflow integration.", BLUE),
        ("~6,300-word whitepaper", "Elegant.md ships in the repo.\nFor stats-literate reviewers.\nIncludes the coverage proof.", AMBER),
    ]
    for title_text, body_text, accent in cards:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, col_h)
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = BG_CARD
        card.line.color.rgb = accent
        card.line.width = Pt(1)
        # Inner text
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        # Title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = title_text
        r.font.size = Pt(20)
        r.font.color.rgb = accent
        r.font.bold = True
        r.font.name = FONT
        # Body
        for line in body_text.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(4)
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(14)
            r2.font.color.rgb = TEXT_PRIMARY
            r2.font.name = FONT
        x = x + col_w + gap
    add_footer(s, 11)
    notes(
        s,
        "FlightCast is MIT-licensed. Five reusable MariaDB-only SQL queries. "
        "Three open issues for the next contributor. The temporal-tables × "
        "conformal-prediction intersection is now public infrastructure for "
        "the MariaDB ecosystem.",
    )


def slide_12_closing(prs):
    s = new_slide(prs)
    add_bg(s)
    add_text(
        s,
        Inches(1),
        Inches(1.6),
        Inches(11.33),
        Inches(1.2),
        "Thank you.",
        size=64,
        color=TEXT_PRIMARY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(2.7),
        Inches(11.33),
        Inches(0.5),
        "Low Yan Cheng  ·  TP070056  ·  APU Malaysia",
        size=18,
        color=VIOLET_SOFT,
        align=PP_ALIGN.CENTER,
    )
    # URL block
    urls_y = Inches(3.8)
    add_text(
        s,
        Inches(1),
        urls_y,
        Inches(11.33),
        Inches(0.4),
        "Live demo  →  github.com/imycc1221/flightcast",
        size=18,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(4.3),
        Inches(11.33),
        Inches(0.4),
        "Hackathon submission  →  github.com/MariaDB-Hackathon-MY-2026/flightcast",
        size=18,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(4.8),
        Inches(11.33),
        Inches(0.4),
        "Whitepaper  →  Elegant.md (in repo)",
        size=18,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(5.3),
        Inches(11.33),
        Inches(0.4),
        "Judge guide  →  docs/JUDGES_TESTING_GUIDE.md  (4 time-budget options)",
        size=18,
        color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(1),
        Inches(6.6),
        Inches(11.33),
        Inches(0.3),
        "Questions welcome.",
        size=14,
        color=TEXT_SECONDARY,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, 12)
    notes(
        s,
        "Thank you for watching. The repo is open, the whitepaper is in the "
        "repo, and the live demo runs locally with one Docker command. "
        "I'm happy to answer questions.",
    )


# ─── Backup Q&A slides ──────────────────────────────────────────────


def slide_backup_a(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "BACKUP  ·  Q&A", color=TEXT_SECONDARY)
    add_title(s, "How realistic is the synthetic demand?", size=34)
    add_subtitle(s, "Tier-aware log-normal noise on real OpenFlights routes.")
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(3.5),
        [
            "Source: OpenFlights public dataset — 50 sampled real airline routes.",
            "Three traffic tiers (hub / mid / thin) with different noise sigmas.",
            "730 days × 50 routes = ~36,500 training rows, log-normal noise.",
            "Drift simulation: σ = 0.10 → σ = 0.22 documented as reproducible recipe in src/flightcast/synth_demand.py.",
        ],
        size=15,
        bullet=BLUE,
    )


def slide_backup_b(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "BACKUP  ·  Q&A", color=TEXT_SECONDARY)
    add_title(s, "Will it scale?", size=34)
    add_subtitle(s, "Performance benchmarks (FOR SYSTEM_TIME AS OF query latency).")
    add_bullets(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(3.5),
        [
            "1 K rows: sub-millisecond temporal lookup.",
            "10 K rows: ~3 ms.",
            "100 K rows: ~25 ms (still under one frame).",
            "Index strategy: idx_airports_iata yielded a 135× speedup on the route-sampling endpoint.",
            "Connection pool (size 8) + HTTP cache headers eliminate the dashboard wedge.",
        ],
        size=15,
        bullet=BLUE,
    )


def slide_backup_c(prs):
    s = new_slide(prs)
    add_bg(s)
    add_eyebrow(s, "BACKUP  ·  Q&A", color=TEXT_SECONDARY)
    add_title(s, "Show me the actual SQL.", size=34)
    add_subtitle(s, "Five hero queries — every one fails on MySQL or PostgreSQL.")
    items = [
        ("1.  Time travel", "FOR SYSTEM_TIME AS OF a date — what did the model predict on this day?"),
        ("2.  Prediction diff", "Two AS OF queries joined — how did forecasts change between v1 and v2?"),
        ("3.  Full audit log", "FOR SYSTEM_TIME ALL — every version of every row, ever."),
        ("4.  Calibration drift", "GROUP BY forecast_run_id — which batches broke?"),
        ("5.  Great-circle distance", "ST_DISTANCE_SPHERE — feature engineering inside the database."),
    ]
    box = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    for i, (title_t, desc_t) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = title_t + "  "
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = VIOLET
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = desc_t
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_PRIMARY
        r2.font.name = FONT


# ─── Main ───────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_insight,
        slide_04_solution,
        slide_05_time_travel,
        slide_06_rainbow,
        slide_07_drift,
        slide_08_math,
        slide_09_architecture,
        slide_10_comparison,
        slide_11_open,
        slide_12_closing,
        slide_backup_a,
        slide_backup_b,
        slide_backup_c,
    ]
    for fn in builders:
        fn(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[OK] Saved {OUT.relative_to(ROOT)} ({len(builders)} slides)")
    # Image audit
    for k, p in IMG.items():
        status = "[OK] " if Path(p).exists() else "[--] "
        print(f"  {status} {k}: {p.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
